
# -*- coding: utf-8 -*-
"""
Generic outline parser for PDF/DOCX/PPTX producing a normalized tree:

{
  "title": "ROOT",
  "children": [ { "title": "...", "children": [...] }, ... ],
  "meta": {...}
}
"""
from __future__ import annotations
from dataclasses import dataclass, field
from typing import List, Dict, Any, Optional, Tuple
import re

# ---- PDF readers (robust import) ----
try:
    from pypdf import PdfReader as PdfReaderNew
except Exception:
    PdfReaderNew = None

try:
    from PyPDF2 import PdfReader as PdfReaderOld
except Exception:
    PdfReaderOld = None

# ---- DOCX / PPTX ----
from docx import Document
from pptx import Presentation

HEADING_HINT = re.compile(
    r"^("
    r"(?:\d+(?:\.\d+){0,3}[\)\.]?\s+)"     # 1 / 1.2 / 1.2.3
    r"|(?:第[一二三四五六七八九十百千]+[章部節項])"  # 第1章/第一章 など
    r"|(?:[A-Z]\.)"                        # A. B. C.
    r"|(?:Appendix|Annex|Chapter|Section)" # 英語見出し語
    r")"
)

BULLET_HINT = re.compile(r"^[\-\u2022\u25CF\u25A0\u25E6\*]\s+")

@dataclass
class Node:
    title: str
    level: int = 1
    children: List["Node"] = field(default_factory=list)
    meta: Dict[str, Any] = field(default_factory=dict)

    def to_dict(self) -> Dict[str, Any]:
        return {"title": self.title, "level": self.level, "children": [c.to_dict() for c in self.children], "meta": self.meta}

def _nest_by_levels(items: List[Tuple[int, str, Dict[str, Any]]], root_title="ROOT") -> Dict[str, Any]:
    root = Node(root_title, level=0)
    stack = [root]
    for lvl, title, meta in items:
        title = re.sub(r"\s+", " ", title.strip())
        if not title:
            continue
        node = Node(title=title, level=max(1, int(lvl)), meta=meta)
        # climb up until parent is above this level
        while stack and node.level <= stack[-1].level:
            stack.pop()
        stack[-1].children.append(node)
        stack.append(node)
    return root.to_dict()

def _guess_level(line: str) -> int:
    # 1.2.3 -> 3, "1) " -> 1, bullets -> 2, otherwise 2
    line = line.strip()
    m = re.match(r"^(\d+(?:\.\d+)+)", line)
    if m:
        return m.group(1).count(".") + 1
    if re.match(r"^\d+[\.\)]", line):
        return 1
    if re.match(r"^[A-Z]\.", line):
        return 1
    if re.match(r"^第[一二三四五六七八九十百千]+[章部節項]", line):
        return 1
    if BULLET_HINT.search(line):
        return 2
    return 2

# -------- PDF --------
def parse_pdf(path: str) -> Dict[str, Any]:
    reader = None
    outline = None
    items: List[Tuple[int, str, Dict[str, Any]]] = []

    # try pypdf first
    if PdfReaderNew is not None:
        try:
            reader = PdfReaderNew(path)
            try:
                outline = reader.outline
            except Exception:
                try:
                    outline = reader.outlines
                except Exception:
                    outline = None
        except Exception:
            reader = None

    # fallback to PyPDF2
    if reader is None and PdfReaderOld is not None:
        try:
            reader = PdfReaderOld(path)
            try:
                outline = reader.outline
            except Exception:
                try:
                    outline = reader.outlines
                except Exception:
                    outline = None
        except Exception:
            reader = None

    def walk(out, level=1):
        for o in out:
            if isinstance(o, list):
                walk(o, level+1)
            else:
                title = str(getattr(o, "title", o))
                items.append((level, title, {"type": "pdf_outline"}))

    # extract outlines if present
    if outline:
        try:
            walk(outline, 1)
        except Exception:
            items = []

    # fallback: heading-like text from first N pages
    if not items and reader is not None:
        pages = min(40, len(reader.pages))
        for i in range(pages):
            try:
                txt = reader.pages[i].extract_text() or ""
            except Exception:
                txt = ""
            for raw in txt.splitlines():
                line = raw.strip()
                if not line:
                    continue
                if HEADING_HINT.search(line) or len(line) <= 32:
                    level = _guess_level(line)
                    items.append((level, line, {"type": "pdf_text", "page": i+1}))

    if not items:
        items = [(1, "Document", {"type": "pdf_empty"})]
    return _nest_by_levels(items, root_title="PDF")

# -------- DOCX --------
def parse_docx(path: str) -> Dict[str, Any]:
    doc = Document(path)
    items: List[Tuple[int, str, Dict[str, Any]]] = []

    for p in doc.paragraphs:
        text = (p.text or "").strip()
        if not text:
            continue
        style = getattr(p.style, "name", "") or ""
        lvl = None

        # Heading styles (both English & Japanese)
        if style.lower().startswith("heading") or "見出し" in style:
            m = re.search(r"(\d+)$", style)
            lvl = int(m.group(1)) if m else 1

        # Numbered headings heuristics
        if lvl is None and HEADING_HINT.search(text):
            m = re.match(r"^(\d+(?:\.\d+)+)", text)
            lvl = m.group(1).count(".") + 1 if m else 2

        if lvl is not None:
            items.append((min(lvl, 6), text, {"type": "docx_heading", "style": style}))

    if not items:
        # fallback: take first non-empty paragraphs as level-1
        for p in doc.paragraphs[:50]:
            t = (p.text or "").strip()
            if t:
                items.append((1, t[:60], {"type": "docx_text"}))

    return _nest_by_levels(items, root_title="DOCX")

# -------- PPTX --------
def parse_pptx(path: str) -> Dict[str, Any]:
    prs = Presentation(path)
    items: List[Tuple[int, str, Dict[str, Any]]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = None
        # title placeholder
        if slide.shapes.title and hasattr(slide.shapes.title, "text"):
            title = slide.shapes.title.text.strip() or None
        # otherwise first text box
        if title is None:
            for shp in slide.shapes:
                if hasattr(shp, "text_frame") and shp.text_frame and shp.text_frame.text.strip():
                    title = shp.text_frame.text.strip().splitlines()[0]
                    break
        items.append((1, f"Slide {idx}: {title or 'Untitled'}", {"type": "pptx_title", "slide": idx}))

        # bullets as level 2/3 depending on indent level
        for shp in slide.shapes:
            if not hasattr(shp, "text_frame") or not shp.text_frame:
                continue
            for p in shp.text_frame.paragraphs:
                txt = "".join(run.text for run in p.runs).strip() or p.text.strip()
                if not txt:
                    continue
                if txt == title:
                    continue
                level = (p.level or 0) + 2  # level 2..N
                items.append((level, txt, {"type": "pptx_bullet", "slide": idx, "indent": p.level or 0}))

    return _nest_by_levels(items, root_title="PPTX")

# -------- Public API --------
def parse_any(path: str) -> Dict[str, Any]:
    low = path.lower()
    if low.endswith(".pdf"):
        return parse_pdf(path)
    if low.endswith(".docx"):
        return parse_docx(path)
    if low.endswith(".pptx"):
        return parse_pptx(path)
    raise ValueError("Unsupported file type")

def flatten_to_depth(tree: Dict[str, Any], max_depth: int) -> Dict[str, Any]:
    """Trim children beyond max_depth (root level=0)."""
    def rec(node: Dict[str, Any], depth: int) -> Dict[str, Any]:
        out = {"title": node.get("title", ""), "level": node.get("level", depth), "meta": node.get("meta", {})}
        if depth >= max_depth:
            out["children"] = []
        else:
            out["children"] = [rec(c, depth+1) for c in node.get("children", [])]
        return out
    return rec(tree, 0)

def collapse_single_chains(tree: Dict[str, Any]) -> Dict[str, Any]:
    """Collapse nodes with a single child into a combined label (for compact mindmap)."""
    def rec(node):
        label = node.get("title", "")
        children = node.get("children", [])
        while len(children) == 1 and children[0].get("children"):
            label = f"{label} / {children[0]['title']}"
            node = children[0]
            children = node.get("children", [])
        return {"title": label, "children": [rec(c) for c in children], "meta": node.get("meta", {})}
    return rec(tree)
